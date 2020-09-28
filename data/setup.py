import argparse

parser = argparse.ArgumentParser()
setup_config = parser.add_argument_group('dataset setting')
setup_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp030')
setup_config.add_argument('--data_name', type=str, dest='data_name', default='gdstd', help='npy name')
setup_config.add_argument('--train_xlsx', type=str, dest='train_xlsx',
                          default='/data/KUGH/mastoid/info/dataset/gdstd_img_190531.xlsx')
setup_config.add_argument('--val_xlsx', type=str, dest='val_xlsx',
                          default='/data/KUGH/mastoid/info/dataset/gdstd_img_190531.xlsx')
setup_config.add_argument('--view_type', type=int, dest='view_type', default=0)
setup_config.add_argument('--image_prop', type=float, dest='image_prop', default=1.5)
setup_config.add_argument('--batch_size', type=int, dest='batch_size', default=60)
setup_config.add_argument('--cuda_device', type=str, dest='cuda_device', default='0')

parser.print_help()
config, unparsed = parser.parse_known_args()

import tensorflow as tf
import numpy as np
import os, re

import pptx
from pptx.util import Inches

import pandas as pd
import pydicom as dicom
import skimage.transform, scipy.misc
import matplotlib.pyplot as plt

import warnings

import sys, os
sys.path.append('/home/chzze/bitbucket/Mastoid_rls')

warnings.filterwarnings('ignore')

DATA_PATH = '/data/KUGH/mastoid/'
INFO_PATH = os.path.join(DATA_PATH, 'info', 'dataset')

IMAGE_WIDTH = 256
IMAGE_HEIGHT = int(IMAGE_WIDTH * float(config.image_prop))
IMAGE_SIZE = [IMAGE_HEIGHT, IMAGE_WIDTH]


def str_extract(string, text):
    extract = re.search(string, text)
    if extract is None:
        matching = None
    else:
        matching = extract.group()
    return matching


def data_setting(npy_name):
    train_xlsx = os.path.join(INFO_PATH, config.train_xlsx)
    val_xlsx = os.path.join(INFO_PATH, config.val_xlsx)

    only_val = train_xlsx == val_xlsx
    npy_path = os.path.join(DATA_PATH, config.exp_name, 'npy')
    data_dir = os.path.join(npy_path, npy_name)
    xlsx_path = os.path.join(DATA_PATH, config.exp_name, 'xlsx')
    if not os.path.exists(xlsx_path): os.makedirs(xlsx_path)
    xlsx = pd.DataFrame()

    xlsx_test = pd.read_excel(val_xlsx)
    xlsx_test['DATA_TYPE'] = 'val'
    xlsx = xlsx.append(xlsx_test, ignore_index=True)

    if only_val is False:
        xlsx_train = pd.read_excel(train_xlsx)
        xlsx_train['DATA_TYPE'] = 'train'
        xlsx = xlsx.append(xlsx_train, ignore_index=True)

    xlsx = xlsx.drop(index=xlsx[np.isnan(xlsx['SPACING0_X'])].index)

    xlsx = xlsx.rename({'REVERSE_AP': 'REVERSE0', 'REVERSE_RT': 'REVERSE1', 'REVERSE_LT': 'REVERSE2'},
                         axis='columns')

    select_col = ['FOLDER_ID', 'PATIENT_ID', 'SERIES_DATE', 'RT_LABEL', 'LT_LABEL',
                  'REVERSE0', 'REVERSE1', 'REVERSE2',
                  'FILENAME_0', 'FILENAME_1', 'FILENAME_2',
                  'IMG_X0', 'IMG_Y0', 'SPACING0_X', 'SPACING0_Y',
                  'IMG_X1', 'IMG_Y1', 'SPACING1_X', 'SPACING1_Y',
                  'IMG_X2', 'IMG_Y2', 'SPACING2_X', 'SPACING2_Y',
                  'TOTAL_LABEL', 'DATA_TYPE']

    xlsx = xlsx[select_col]

    if not os.path.isfile(os.path.join(xlsx_path, config.data_name+'_info.xlsx')):

        xlsx_side = pd.DataFrame({'FOLDER_ID': pd.Series(), 'PATIENT_ID': pd.Series(),
                                  'SERIES_DATE': pd.Series(), 'DATA_TYPE': pd.Series(),
                                  'FILENAME_0': pd.Series(), 'FILENAME_1': pd.Series(),
                                  'SIDE': pd.Series(), 'LABEL': pd.Series(),
                                  'REVERSE0': pd.Series(), 'REVERSE1': pd.Series(),
                                  'IMG_X0': pd.Series(), 'IMG_Y0': pd.Series(),
                                  'SPACING0_X': pd.Series(), 'SPACING0_Y': pd.Series(),
                                  'IMG_X1': pd.Series(), 'IMG_Y1': pd.Series(),
                                  'SPACING1_X': pd.Series(), 'SPACING1_Y': pd.Series()
                                  })

        sides_match = {'RT': '1', 'LT': '2'}

        j = 0
        for idx in xlsx.index:
            for side in list(sides_match.keys()):
                lat_idx = sides_match[side]

                xlsx_side.loc[j, 'SIDE'] = side
                xlsx_side.loc[j, ['FOLDER_ID', 'PATIENT_ID']] = xlsx.loc[idx, ['FOLDER_ID', 'PATIENT_ID']]
                xlsx_side.loc[j, ['SERIES_DATE', 'DATA_TYPE']] = xlsx.loc[idx, ['SERIES_DATE', 'DATA_TYPE']]

                xlsx_side.loc[j, 'LABEL'] = xlsx.loc[idx, side+'_LABEL']
                xlsx_side.loc[j, 'FILENAME_0'] = xlsx.loc[idx, 'FILENAME_0']
                xlsx_side.loc[j, ['IMG_X0', 'IMG_Y0']] = xlsx.loc[idx, ['IMG_X0', 'IMG_Y0']]
                xlsx_side.loc[j, ['SPACING0_X', 'SPACING0_Y']] = xlsx.loc[idx, ['SPACING0_X', 'SPACING0_Y']]
                xlsx_side.loc[j, 'REVERSE0'] = xlsx.loc[idx, 'REVERSE0']

                xlsx_side.loc[j, 'FILENAME_1'] = xlsx.loc[idx, 'FILENAME_' + lat_idx]
                xlsx_side.loc[j, 'IMG_X1'] = xlsx.loc[idx, 'IMG_X' + lat_idx]
                xlsx_side.loc[j, 'IMG_Y1'] = xlsx.loc[idx, 'IMG_Y' + lat_idx]
                xlsx_side.loc[j, 'SPACING1_X'] = xlsx.loc[idx, 'SPACING' + lat_idx + '_X']
                xlsx_side.loc[j, 'SPACING1_Y'] = xlsx.loc[idx, 'SPACING' + lat_idx + '_Y']

                xlsx_side.loc[j, 'REVERSE1'] = xlsx.loc[idx, 'REVERSE'+lat_idx]

                j += 1

        xlsx_side.to_excel(os.path.join(xlsx_path, config.data_name+'_info.xlsx'))
    else:
        xlsx_side = pd.read_excel(os.path.join(xlsx_path, config.data_name+'_info.xlsx'))

    def calculate_crop_range_ap(height, width, spacing_x, spacing_y, radius_w, side):
        center_y = int(height * 0.6)
        if side == 'RT':
            center_x = int(width / 4)
        elif side == 'LT':
            center_x = 3*int(width / 4)
        else:
            raise ValueError('Invalid side!')

        radius_h = int(radius_w * float(config.image_prop))
        x1, y1 = int(center_x - radius_w / spacing_x), int(center_y - radius_h / spacing_y)
        x2, y2 = int(center_x + radius_w / spacing_x), int(center_y + radius_h / spacing_y)

        return [x1, y1, x2, y2]

    def calculate_crop_range_lat(height, width, spacing_x, spacing_y, radius):
        center_x, center_y = int(width / 2), int(height / 2)

        x1, y1 = int(center_x - radius / spacing_x), int(center_y - radius / spacing_y)
        x2, y2 = int(center_x + radius / spacing_x), int(center_y + radius / spacing_y)

        return [x1, y1, x2, y2]

    def concatenate_id(folder_id, side):
        return '-'.join([folder_id, side])

    def add_value(xlsx, view_type, data_type):
        view_name = 'FILENAME_' + view_type
        x, y = 'IMG_X' + view_type, 'IMG_Y' + view_type
        sx, sy = 'SPACING'+view_type+'_X', 'SPACING'+view_type+'_Y'

        xlsx['FILENAMES'+view_type] = xlsx.apply(
            lambda row: os.path.join(DATA_PATH, row['FOLDER_ID'], row[view_name]), axis=1)

        if view_type == '0':
            xlsx[view_name + '_INFO'] = xlsx.apply(
                lambda row: calculate_crop_range_ap(height=row[y], width=row[x],
                                                    spacing_x=row[sx], spacing_y=row[sy],
                                                    radius_w=60, side=row['SIDE']), axis=1)
        else:
            xlsx[view_name + '_INFO'] = xlsx.apply(
                lambda row: calculate_crop_range_lat(height=row[y], width=row[x],
                                                     spacing_x=row[sx], spacing_y=row[sy],
                                                     radius=70), axis=1)

        xlsx['LABELS' + view_type] = xlsx.apply(lambda row:
                                                row[view_name+'_INFO']+
                                                [int(row['REVERSE'+view_type])]+
                                                [0 if x == 0 else 1 for x in [row['LABEL']]], axis=1)

        xlsx['FOLDER_SIDE'] = xlsx.apply(
            lambda row: concatenate_id(folder_id=row['FOLDER_ID'], side=row['SIDE']), axis=1)

        files_view = xlsx[xlsx['DATA_TYPE'] == data_type]['FILENAMES' + view_type].values
        labels_view = xlsx[xlsx['DATA_TYPE'] == data_type]['LABELS' + view_type].values
        names_view = xlsx[xlsx['DATA_TYPE'] == data_type]['FOLDER_SIDE'].values

        return files_view, labels_view, names_view

    def csv_files_labels(data_type):
        file0, label0, names = add_value(xlsx_side, '0', data_type)
        file1, label1, _ = add_value(xlsx_side, '1', data_type)

        csv_info = pd.DataFrame({'FILENAMES0': pd.Series(file0), 'LABELS0': pd.Series(label0),
                                 'FILENAMES1': pd.Series(file1), 'LABELS1': pd.Series(label1),
                                 'ID': pd.Series(names)
                                })
        csv_info.to_csv(os.path.join(xlsx_path, npy_name+'_'+data_type+'.csv'))
        return file0, file1, label0, label1, names

    file0_val, file1_val, label0_val, label1_val, id_val = csv_files_labels('val')

    if only_val:
        dataset = DataSettingV1(data_dir=data_dir, batch_size=config.batch_size, only_val=only_val,
                                file0_val=file0_val, label0_val=label0_val,
                                file1_val=file1_val, label1_val=label1_val,
                                id_val=id_val
                                )
    else:
        file0_train, file1_train, label0_train, label1_train, id_train = csv_files_labels('train')

        dataset = DataSettingV1(data_dir=data_dir, batch_size=config.batch_size, only_val=only_val,
                                file0_val=file0_val, label0_val=label0_val,
                                file1_val=file1_val, label1_val=label1_val,
                                id_val=id_val,
                                file0_train=file0_train, label0_train=label0_train,
                                file1_train=file1_train, label1_train=label1_train,
                                id_train=id_train
                                )
    return dataset


class DataSettingV1():
    def __init__(self, data_dir, batch_size, only_val, **kwargs):
        if not os.path.exists(data_dir):
            if 'id_train' in kwargs:
                train_x0, train_y0 = kwargs['file0_train'], kwargs['label0_train']
                train_x1, train_y1 = kwargs['file1_train'], kwargs['label1_train']
                train_z = kwargs['id_train']

            if 'id_val' in kwargs:
                val_x0, val_y0 = kwargs['file0_val'], kwargs['label0_val']
                val_x1, val_y1 = kwargs['file1_val'], kwargs['label1_val']
                val_z = kwargs['id_val']

            else:
                raise AssertionError('images or labels must be provided. please check npy file.')

            data_root = os.path.split(data_dir)[0]
            if not os.path.exists(data_root): os.makedirs(data_root)

            if only_val:
                np.save(data_dir, {'val_x0': val_x0, 'val_y0': val_y0,
                                   'val_x1': val_x1, 'val_y1': val_y1, 'val_z': val_z
                                   })

            else:
                np.save(data_dir, {'train_x0': train_x0, 'train_y0': train_y0,
                                   'train_x1': train_x1, 'train_y1': train_y1, 'train_z': train_z,
                                   'val_x0': val_x0, 'val_y0': val_y0,
                                   'val_x1': val_x1, 'val_y1': val_y1, 'val_z': val_z
                                   })
        else:
            pre_built = np.load(data_dir).item()

            if only_val:
                val_x0, val_y0 = pre_built['val_x0'], pre_built['val_y0']
                val_x1, val_y1 = pre_built['val_x1'], pre_built['val_y1']
                val_z = pre_built['val_z']
                self.data_length = len(val_z)

            else:
                train_x0, train_y0 = pre_built['train_x0'], pre_built['train_y0']
                train_x1, train_y1 = pre_built['train_x1'], pre_built['train_y1']
                train_z = pre_built['train_z']

                val_x0, val_y0 = pre_built['val_x0'], pre_built['val_y0']
                val_x1, val_y1 = pre_built['val_x1'], pre_built['val_y1']
                val_z = pre_built['val_z']

                self.data_length = len(train_z) + len(val_z)

        self.val = self.SubDataSetting((val_x0, val_y0, val_x1, val_y1, val_z),
                                       batch_size=batch_size, shuffle=False, augmentation=False)
        self.id_val = val_z

        if only_val is False:
            index0 = np.asarray([v[-1] for v in train_y0])
            index1 = np.asarray([v[-1] for v in train_y1])

            train_x0 = np.concatenate([train_x0, train_x0[index0 == 0]], axis=0)
            train_y0 = np.concatenate([train_y0, train_y0[index0 == 0]], axis=0)

            train_x1 = np.concatenate([train_x1, train_x1[index1 == 0]], axis=0)
            train_y1 = np.concatenate([train_y1, train_y1[index1 == 0]], axis=0)

            train_z = np.concatenate([train_z, train_z[index0 == 0]], axis=0)

            np.random.seed(20190617)
            p = np.random.permutation(len(train_x0))

            train_x0, train_y0 = train_x0[p], train_y0[p]
            train_x1, train_y1 = train_x1[p], train_y1[p]
            train_z = train_z[p]

            self.train = self.SubDataSetting((train_x0, train_y0, train_x1, train_y1, train_z),
                                             batch_size=batch_size, shuffle=True, augmentation=True)
            self.id_train  = train_z

    class SubDataSetting():
        def __init__(self, filenames_n_labels, num_epochs=1, batch_size=1, shuffle=False, augmentation=False):

            self.file0, self.label0, self.file1, self.label1, self.id = filenames_n_labels
            self.data_length = len(self.id)

            dataset = tf.data.Dataset.from_tensor_slices(tensors=
                                                         (self.file0, [v for v in self.label0],
                                                          self.file1, [v for v in self.label1],
                                                          [v for v in self.id]
                                                          ))
            if shuffle:
                dataset = dataset.shuffle(buffer_size=batch_size * 100, reshuffle_each_iteration=True)

            def is_reverse(img, bw, reverse):
                img_rows, img_cols = img.shape[0], img.shape[1]
                img_patch = img[0:int(0.1 * img_rows), 0:int(0.1 * img_cols)]
                patch_mean = np.mean(img_patch)
                img_min, img_max = np.min(img), np.max(img)
                min_diff, max_diff = patch_mean - img_min, img_max - patch_mean
                if bw == 'black':
                    return min_diff >= max_diff
                elif bw == 'white':
                    if reverse == 0:
                        return min_diff <= max_diff
                    else:
                        return min_diff >= max_diff
                else:
                    return None

            def dcm_read_by_ftn(file0, label0, file1, label1, name, augmentation):
                def each_read(filename, label, bw, augmentation):
                    dcm_info = dicom.read_file(filename.decode())
                    x1, y1, x2, y2 = label[:-2]
                    reverse = label[-2]

                    if augmentation:
                        shift_x = np.random.randint(dcm_info.Columns // 20)
                        shift_y = np.random.randint(dcm_info.Rows // 20)

                        shift_x = -shift_x if np.random.rand() <= 0.5 else shift_x
                        shift_y = -shift_y if np.random.rand() <= 0.5 else shift_y

                        x1, y1 = x1-shift_x, y1-shift_y
                        x2, y2 = x2-shift_x, y2-shift_y
                    image = dcm_info.pixel_array

                    if len(image.shape) >= 3:
                        image = image[:, :, 0]

                    if is_reverse(image, bw, reverse):
                        white_image = np.full_like(image, np.max(image), image.dtype)
                        image = np.subtract(white_image, image)

                    image = image[max(0, y1):min(dcm_info.Rows, y2), max(0, x1):min(dcm_info.Columns, x2)]

                    if bw == 'black':
                        image_size = [IMAGE_HEIGHT, IMAGE_WIDTH]
                    elif bw == 'white':
                        image_size = [IMAGE_WIDTH, IMAGE_WIDTH]
                    else:
                        raise ValueError('Error! Invalid bw type.')

                    image = np.expand_dims(skimage.transform.resize(image, image_size, preserve_range=True), axis=-1)

                    if augmentation and np.random.randint(2) == 1:
                        image = np.fliplr(image)

                    image = (image - np.mean(image)) / np.std(image)

                    return image.astype(np.float32), np.int64(label[-1])

                f0, l0 = each_read(file0, label0.astype(int), 'black', augmentation)
                f1, _ = each_read(file1, label1.astype(int), 'white', augmentation)
                nm = name.decode()

                return f0, f1, l0, nm

            dataset = dataset.map(num_parallel_calls=8,
                                  map_func=lambda file0, label0, file1, label1, name:
                                  tuple(tf.py_func(func=dcm_read_by_ftn,
                                                   inp=[file0, label0, file1, label1, name, augmentation],
                                                   Tout=[tf.float32, tf.float32, tf.int64, tf.string]
                                                   )))
            if num_epochs == 0:
                dataset = dataset.repeat(count=num_epochs)  # raise out-of-range error when num_epochs done
                dataset = dataset.batch(batch_size=batch_size, drop_remainder=True)
            else:
                dataset = dataset.batch(batch_size)
                dataset = dataset.repeat(count=num_epochs)
            iterator = dataset.make_initializable_iterator()

            self.init_op = iterator.initializer
            self.next_batch = iterator.get_next()


def test_dataset_loader(dataset):
    check_path = os.path.join(DATA_PATH, config.exp_name, 'view')
    if not os.path.exists(check_path):
        os.makedirs(check_path)

    with tf.Session() as sess:

        sess.run(dataset.val.init_op)
        num_examples, next_batch = dataset.val.data_length, dataset.val.next_batch

        count = 0
        inch_w, inch_h = 10, 6
        # batch_size = inch_w * inch_h * 2 * 2
        batch_size = config.batch_size
        num_iter = int(np.ceil(float(num_examples) / batch_size))
        print('num_iter: ', num_iter)

        prs = pptx.Presentation()
        prs.slide_width = Inches(inch_w * 2)
        prs.slide_height = Inches(inch_h * 2)

        while count < num_iter:
            img0, img1, lbl, name = sess.run(next_batch)

            if config.view_type == 0:
                show_img = img0
            elif config.view_type == 1:
                show_img = img1
            else:
                show_img = None

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            show_images(show_img, name, num_rows=inch_h, num_cols=inch_w, figsize=(inch_w*2, inch_h*2))
            fig_name = '_'.join([config.exp_name, config.data_name,
                                 'type%d' % config.view_type, '%03d' % count])+'.png'
            fig_path = os.path.join(check_path, fig_name)
            plt.savefig(fig_path, bbox_inches='tight', dpi=100, pil_kwargs={"compression": "png_lzw"})

            slide.shapes.add_picture(fig_path, Inches(0), Inches(0),
                                     width=Inches(inch_w*2), height=Inches(inch_h*2))
            os.remove(fig_path)
            count += 1

            if count % 10 == 0:
                print(count)

    ppt_name = os.path.join(check_path, '_'.join([config.exp_name, config.data_name,
                                                  'type%d' % config.view_type])+'.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_images(images, names, num_rows=6, num_cols=10, figsize=(10*2, 6*2)):
    plt.figure(figsize=figsize)
    num_figs = images.shape[0]  # num_rows * num_cols

    for j in range(num_figs):
        plt.subplot(num_rows, num_cols, j + 1)
        plt.imshow(np.squeeze(images[j]), cmap='gray')
        plt.axis('off')


if __name__ == '__main__':
    os.environ['CUDA_VISIBLE_DEVICES'] = config.cuda_device
    d_set = data_setting(npy_name='_'.join([config.exp_name, config.data_name]) + '.npy')
    test_dataset_loader(d_set)
