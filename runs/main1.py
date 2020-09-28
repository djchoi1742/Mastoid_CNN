import argparse

license = """
Copyright ⓒ Dongjun Choi, Kyong Joon Lee
Department of Radiology at Seoul National University Bundang Hospital. \n
If you have any question, please email us for assistance: chzze4582@gmail.com \n """
parser = argparse.ArgumentParser(formatter_class=argparse.RawDescriptionHelpFormatter, \
                                 description='', epilog=license, add_help=False)

network_config = parser.add_argument_group('network setting (must be provided)')

network_config.add_argument('--data_path', type=str, dest='data_path', default='/data/KUGH/mastoid/')
network_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp030')
network_config.add_argument('--model_name', type=str, dest='model_name', default='Model51')
network_config.add_argument('--train', type=lambda x: x.title() in str(True), dest='train', default=False)
network_config.add_argument('--batch_size', type=int, dest='batch_size', default=12)
network_config.add_argument('--numEpoch', type=int, dest='num_epoch', default=0)  # infinite loop
network_config.add_argument('--trial_serial', type=int, dest='trial_serial', default=1)
network_config.add_argument('--npy_name', type=str, dest='npy_name', default='exp030_trval_all.npy')
network_config.add_argument('--max_keep', type=int, dest='max_keep', default=10)  # only use training
network_config.add_argument('--num_weight', type=int, dest='num_weight', default=20)  # only use validation

parser.print_help()
config, unparsed = parser.parse_known_args()

import sys, os
sys.path.append('/home/chzze/bitbucket/Mastoid_rls')

os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'  # ignore SSE instruction warning on tensorflow

import tensorflow as tf
import numpy as np
import sklearn.metrics  # roc curve
import matplotlib.pyplot as plt
import pandas as pd
import json, pptx, re
from pptx.util import Inches


trial_serial_str = '%03d' % (config.trial_serial)
log_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'logs-%s' % trial_serial_str)
result_path = os.path.join(config.data_path, config.exp_name, config.model_name, 'result-%s' % trial_serial_str)
ckpt_path = os.path.join(result_path, 'ckpt')
npy_path = os.path.join(config.data_path, config.exp_name, 'npy')

cam_path = os.path.join(config.data_path, 'cam')
ppt_path = os.path.join(config.data_path, 'pptx')

if not os.path.exists(result_path): os.makedirs(result_path)
from data.setup import DataSettingV1

import models.network as network
from tf_utils.tboard import Tensorboard

print('/'.join([npy_path, config.npy_name]))
dataset = DataSettingV1(data_dir=os.path.join(npy_path, config.npy_name), batch_size=config.batch_size,
                        only_val=bool(1 - config.train))

infer_name = 'Inference' + config.model_name
model = getattr(network, infer_name)(trainable=config.train)

sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
sess_config.gpu_options.allow_growth = True
sess = tf.Session(config=sess_config)


def prep_ckpt_rename(prep_ckpt):
    return re.sub('/workspace/Mastoid/', config.data_path, prep_ckpt)


def training():
    tensorboard = Tensorboard(log_dir=log_path, overwrite=True)
    loss_rec = tf.get_variable(name='Loss', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                               collections=['scalar'])
    auc_rec = tf.get_variable(name='AUC', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                              collections=['scalar'])
    accuracy_rec = tf.get_variable(name='Accuracy', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                              collections=['scalar'])

    tensorboard.init_scalar(collections=['scalar'])
    tensorboard.init_images(collections=['images'], num_outputs=4)

    sess.run(tf.global_variables_initializer())
    saver = tf.train.Saver(max_to_keep=config.max_keep)

    result_name = '_'.join([config.exp_name, config.model_name, trial_serial_str])+'.csv'
    auc_csv = pd.DataFrame({'WEIGHT_PATH': pd.Series(), 'AUC': pd.Series()})

    current_step, current_epoch = None, None
    perf_per_epoch, max_perf_per_epoch, max_current_step = [], [], []
    try:
        while True:
            sess.run([dataset.train.init_op, dataset.val.init_op])
            train_loss_batch, train_acc_batch = [], []
            train_x, train_y = [], []

            train_length = dataset.train.data_length
            num_iter_train = int(np.ceil(float(train_length) / config.batch_size))
            train_step = 0

            feed_dict = {}
            while train_step < num_iter_train:
                images0, images1, labels, names = sess.run(dataset.train.next_batch)

                if config.model_name == 'Model51':
                    images = images0
                    feed_dict = {model.images: images, model.labels: labels, model.is_training: True}

                elif config.model_name == 'Model52':
                    images = images1
                    feed_dict = {model.images: images, model.labels: labels, model.is_training: True}

                else:
                    raise ValueError('Invalid model name!')

                _, train_loss, train_prob, train_acc = \
                    sess.run([model.train, model.loss, model.prob, model.accuracy], feed_dict=feed_dict)

                current_step, current_epoch = sess.run([tf.train.get_global_step(), model.global_epoch])
                sys.stdout.write('Step: {0:>4d} ({1})\r'.format(current_step, current_epoch))

                train_x.extend(train_prob)
                train_y.extend(labels)
                train_acc_batch.append(train_acc)
                train_loss_batch.append(train_loss)

                train_step += 1

            sess.run(tf.assign_add(model.global_epoch, 1))

            fpr, tpr, _ = sklearn.metrics.roc_curve(train_y, train_x, drop_intermediate=False)
            train_auc = sklearn.metrics.auc(fpr, tpr)

            feed_dict.update({loss_rec: np.mean(train_loss_batch), auc_rec: train_auc,
                              accuracy_rec: np.mean(train_acc_batch)})

            tensorboard.add_summary(sess=sess, feed_dict=feed_dict, log_type='train')

            val_length = dataset.val.data_length
            num_iter_val = int(np.ceil(float(val_length) / config.batch_size))
            val_step = 0

            val_loss_batch, val_acc_batch = [], []
            val_x, val_y = [], []

            feed_dict = {}
            while val_step < num_iter_val:
                sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_loss_batch),
                                                                 -(-dataset.val.data_length // config.batch_size)))
                images0, images1, labels, names = sess.run(dataset.val.next_batch)

                if config.model_name == 'Model51':
                    images = images0
                    feed_dict = {model.images: images, model.labels: labels, model.is_training: False}

                elif config.model_name == 'Model52':
                    images = images1
                    feed_dict = {model.images: images, model.labels: labels, model.is_training: False}

                else:
                    raise ValueError('Invalid model name!')

                val_loss, val_prob, val_acc = sess.run([model.loss, model.prob, model.accuracy], feed_dict=feed_dict)

                val_x.extend(val_prob)
                val_y.extend(labels)
                val_acc_batch.append(val_acc)
                val_loss_batch.append(val_loss)

                val_step += 1

            fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, val_x, drop_intermediate=False)
            val_auc = sklearn.metrics.auc(fpr, tpr)

            feed_dict.update({loss_rec: np.mean(val_loss_batch), auc_rec: val_auc,
                              accuracy_rec: np.mean(val_acc_batch)})
            tensorboard.add_summary(sess=sess, feed_dict=feed_dict, log_type='val')
            tensorboard.display_summary(time_stamp=True)

            current_epoch += 1
            if current_epoch % 1 == 0:
                perf_per_epoch.append(val_auc)

                if current_epoch < config.max_keep + 1:
                    max_current_step.append(current_step)
                    max_perf_per_epoch.append(val_auc)

                    saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'),
                               global_step=current_step)
                    auc_csv.loc[current_step, 'WEIGHT_PATH'] = \
                        os.path.join(log_path, 'model.ckpt-'+str(current_step))
                    auc_csv.loc[current_step, 'AUC'] = val_auc

                elif val_auc > min(auc_csv['AUC'].tolist()):
                    auc_csv = auc_csv.drop(max_current_step[0])
                    max_current_step.pop(0)
                    max_current_step.append(current_step)
                    max_perf_per_epoch.pop(0)
                    max_perf_per_epoch.append(val_auc)

                    saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'),
                               global_step=current_step)
                    auc_csv.loc[current_step, 'WEIGHT_PATH'] = \
                        os.path.join(log_path, 'model.ckpt-'+str(current_step))
                    auc_csv.loc[current_step, 'AUC'] = val_auc

                auc_csv.to_csv(os.path.join(result_path, result_name))

                if current_epoch == 100: break
        print('Training Complete...\n')
        sess.close()

    except KeyboardInterrupt:
        print('Result saved')
        auc_csv.to_csv(os.path.join(result_path, result_name))


def validation():
    sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
    sess_config.gpu_options.allow_growth = True

    saver = tf.train.Saver()
    num_examples = len(dataset.id_val)

    ckpt = tf.train.get_checkpoint_state(log_path)
    if not ckpt:
        raise ValueError('No checkpoint found in ' + log_path)

    weight_auc_path = os.path.join(config.data_path, config.exp_name, config.model_name,
                                   'result-%03d' %  config.trial_serial)
    weight_auc_csv = pd.read_csv(os.path.join(weight_auc_path, '_'.join([config.exp_name, config.model_name,
                                                                         '%03d' % config.trial_serial])+'.csv'))
    weight_auc_csv = weight_auc_csv.sort_values('AUC', ascending=False)
    all_ckpt_paths = list(weight_auc_csv['WEIGHT_PATH'][0:int(config.num_weight)])

    num_ckpt = len(all_ckpt_paths)
    print('num_ckpt: ', num_ckpt)

    imgs = np.zeros([num_examples, model.img_h, model.img_w, model.img_c])
    cams = np.zeros([num_ckpt, num_examples, model.img_h, model.img_w, model.img_c])

    lbls = np.zeros([num_examples, ], dtype=np.int32)
    probs = np.zeros([num_ckpt, num_examples, 1])

    val_x, val_y = None, None
    for ckpt_idx, ckpt_path in enumerate(all_ckpt_paths):
        print('Restoring: ' + ckpt_path)

        sess = tf.Session(config=sess_config)
        saver.restore(sess, ckpt_path)

        sess.run(dataset.val.init_op)
        val_x, val_y = [], []

        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        step = 0

        while step < num_iter:
            sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_y) // config.batch_size,
                             -(-dataset.val.data_length // config.batch_size)))

            images0, images1, labels, names = sess.run(dataset.val.next_batch)

            if config.model_name == 'Model51':
                images = images0
                feed_dict = {model.images: images, model.labels: labels, model.is_training: False}

            elif config.model_name == 'Model52':
                images = images1
                feed_dict = {model.images: images, model.labels: labels, model.is_training: False}

            else:
                raise ValueError('Invalid model name!')

            val_loss, val_prob, val_acc = \
                sess.run([model.loss, model.prob, model.accuracy], feed_dict=feed_dict)

            val_x.extend(val_prob)
            val_y.extend(labels)

            cam = sess.run(model.local, feed_dict=feed_dict)

            cams[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = cam
            probs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = val_prob

            if ckpt_idx == 0:
                imgs[step * config.batch_size:step * config.batch_size + len(labels)] = images
                lbls[step * config.batch_size:step * config.batch_size + len(labels)] = labels

            step += 1

        sess.close()

    probs, cams = np.mean(probs, axis=0), np.mean(cams, axis=0)
    id_test = dataset.id_val

    prob_1 = np.squeeze(np.array(probs))

    result_csv = pd.DataFrame({'NUMBER': id_test, 'PROB': prob_1, 'LABEL': np.array(lbls)})
    result_name = '_'.join([config.model_name, config.npy_name, trial_serial_str,
                            '%03d' % (config.num_weight)])+'.csv'
    result_csv.to_csv(os.path.join(result_path, result_name), index=False)

    fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, prob_1, drop_intermediate=False)
    val_auc = sklearn.metrics.auc(fpr, tpr)

    print('Validation AUC: ', val_auc)
    print('Validation Complete...\n')

    prs = pptx.Presentation()
    prs.slide_width, prs.slide_height = Inches(8*2), Inches(5*2)

    plt_batch = 20
    plt_step = 0
    plt_iter, plt_examples = int(np.ceil(num_examples / plt_batch)), num_examples

    while plt_step < plt_iter:

        if plt_examples >= plt_batch:
            len_batch = plt_batch
        else:
            len_batch = plt_examples

        images_batch = imgs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        labels_batch = lbls[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        names_batch = id_test[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        probs_batch = probs[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        cams_batch = cams[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        show_cam(cams_batch, probs_batch, images_batch, labels_batch, names_batch, 'LABEL')
        fig_name = '_'.join([config.model_name, config.npy_name, trial_serial_str,
                             '%03d' % plt_step]) + '.png'
        fig_path = os.path.join(cam_path, fig_name)
        plt.savefig(fig_path, bbox_inches='tight')
        slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(8 * 2))
        os.remove(fig_path)
        plt_step += 1
        plt_examples -= plt_batch

    print('plt_examples check: ', plt_examples)
    ppt_name = os.path.join(ppt_path, '_'.join([config.model_name, config.npy_name, trial_serial_str,
                                                '%03d' % config.num_weight]) + '.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_cam(cams, probs, images, labels, names, side_label, num_rows=5, num_cols=8, figsize=(8*2, 5*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=figsize)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    for i in range(batch_size):
        prob = '%.2f' % probs[i]
        lbl = int(labels[i])
        show_image = np.squeeze(images[i])
        cam = np.squeeze(cams[i])
        img_row, img_col = int(i % num_rows), int(i / num_rows) * 2

        ori_title = ' '.join([names[i], side_label + ': '+str(lbl)])
        cam_title = side_label+' Pred: '+str(prob)

        ax[img_row, img_col].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(cam, cmap=plt.cm.jet, alpha=0.5, interpolation='nearest')

        if (lbl == 0 and probs[i] < 0.5) or (lbl == 1 and probs[i] >= 0.5):
            txt_color = 'blue'
        else:
            txt_color = 'red'
        ax[img_row, img_col].set_title(ori_title, fontsize=7, color=txt_color)
        ax[img_row, img_col+1].set_title(cam_title, fontsize=7, color=txt_color)


if __name__ == '__main__':
    if config.train:
        print('Training')
        training()
    else:
        print('Validation')
        validation()